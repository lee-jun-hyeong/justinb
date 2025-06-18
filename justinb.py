import math
import json
import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
import os
import time
import glob
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# PyInstaller 실행 여부에 따른 경로 처리
def resource_path(relative_path):
    try:
        # PyInstaller로 빌드된 exe에서의 경로
        base_path = sys._MEIPASS
    except Exception:
        # 일반 실행 시
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)

# JSON 파일 읽기
with open(resource_path("KorRV.json"), "r", encoding="utf-8") as f:
    bible_data = json.load(f)

full_names = {
    "창": "창세기",
    "출": "출애굽기",
    "레": "레위기",
    "민": "민수기",
    "신": "신명기",
    "수": "여호수아",
    "삿": "사사기",
    "룻": "룻기",
    "삼상": "사무엘상",
    "삼하": "사무엘하",
    "왕상": "열왕기상",
    "왕하": "열왕기하",
    "대상": "역대상",
    "대하": "역대하",
    "스": "에스라",
    "느": "느헤미야",
    "에": "에스더",
    "욥": "욥기",
    "시": "시편",
    "잠": "잠언",
    "전": "전도서",
    "아": "아가",
    "사": "이사야",
    "렘": "예레미야",
    "애": "예레미야애가",
    "겔": "에스겔",
    "단": "다니엘",
    "호": "호세아",
    "욜": "요엘",
    "암": "아모스",
    "옵": "오바댜",
    "욘": "요나",
    "미": "미가",
    "나": "나훔",
    "합": "하박국",
    "습": "스바냐",
    "학": "학개",
    "슥": "스가랴",
    "말": "말라기",
    "마": "마태복음",
    "막": "마가복음",
    "눅": "누가복음",
    "요": "요한복음",
    "행": "사도행전",
    "롬": "로마서",
    "고전": "고린도전서",
    "고후": "고린도후서",
    "갈": "갈라디아서",
    "엡": "에베소서",
    "빌": "빌립보서",
    "골": "골로새서",
    "살전": "데살로니가전서",
    "살후": "데살로니가후서",
    "딤전": "디모데전서",
    "딤후": "디모데후서",
    "딛": "디도서",
    "몬": "빌레몬서",
    "히": "히브리서",
    "약": "야고보서",
    "벧전": "베드로전서",
    "벧후": "베드로후서",
    "요일": "요한1서",
    "요이": "요한2서",
    "요삼": "요한3서",
    "유": "유다서",
    "계": "요한계시록"
}


def find_book(book_name):
    for book in bible_data["books"]:
        if book["name"] == book_name:
            return book
    print(f"책 '{book_name}'를 찾지 못했어")
    return None

def get_verses(book, chapter_num, verse_start=None, verse_end=None):
    chapter_num = int(chapter_num)
    for chapter in book["chapters"]:
        if int(chapter["chapter"]) == chapter_num:
            verses = chapter.get("verses", [])
            if verse_start is None:
                # 절 번호 포함해서 반환
                return [(chapter_num, int(v["verse"]), v["text"]) for v in verses]
            else:
                verse_start = int(verse_start)
                verse_end = int(verse_end)
                selected = []
                for v in verses:
                    verse_num = int(v["verse"])
                    if verse_start <= verse_num <= verse_end:
                        selected.append((chapter_num, verse_num, v["text"]))
                if not selected:
                    print(f"{chapter_num}장 {verse_start}-{verse_end}절이 없어.")
                return selected
    print(f"장 {chapter_num}를 찾지 못했어.")
    return []


def parse_reference(ref):
    pattern = r'([가-힣]+)(\d*)(?::(\d+)(?:-(\d+))?)?'
    m = re.match(pattern, ref)
    if not m:
        print(f"파싱 실패 똑바로 입력하셈: {ref}")
        return None

    book_abbr = m.group(1)  # 이제 약어를 그대로 책 이름으로 사용
    chapter_str = m.group(2)
    verse_start_str = m.group(3)
    verse_end_str = m.group(4)

    # abbreviations 사용 안 하므로, book_abbr가 JSON 데이터에 있는지 find_book에서 체크
    book = find_book(book_abbr)
    if not book:
        print(f"책 이름 '{book_abbr}'을(를) 찾을 수 없습니다.")
        return None

    chapters = book.get("chapters", [])

    if not chapter_str:
        start_chapter = 1
        start_verse = 1
        end_chapter = int(chapters[-1]["chapter"])
        end_verse = int(chapters[-1]["verses"][-1]["verse"])
        return (book_abbr, start_chapter, start_verse, (end_chapter, end_verse))

    chapter = int(chapter_str)
    if not verse_start_str:
        chapter_data = next((c for c in chapters if int(c["chapter"]) == chapter), None)
        if not chapter_data:
            print(f"{book_abbr} {chapter}장을 찾을 수 없습니다.")
            return None
        verse_list = chapter_data.get("verses", [])
        verse_start = 1
        verse_end = int(verse_list[-1]["verse"]) if verse_list else 1
        return (book_abbr, chapter, verse_start, verse_end)

    verse_start = int(verse_start_str)
    verse_end = int(verse_end_str) if verse_end_str else verse_start
    return (book_abbr, chapter, verse_start, verse_end)


def get_texts_from_reference(ref):
    parsed = parse_reference(ref)
    if not parsed:
        return None
    book_name, chapter_start, verse_start, verse_end = parsed
    book = find_book(book_name)
    if not book:
        print(f"'{book_name}' 책을 찾지 못했습니다.")
        return None

    if isinstance(verse_end, tuple):
        last_chapter, last_verse = verse_end
        texts = []
        for ch in range(chapter_start, last_chapter + 1):
            if ch == chapter_start:
                start_v = verse_start
            else:
                start_v = 1

            if ch == last_chapter:
                end_v = last_verse
            else:
                ch_data = next((c for c in book["chapters"] if int(c["chapter"]) == ch), None)
                if not ch_data:
                    print(f"{book_name} {ch}장을 찾을 수 없습니다.")
                    return None
                end_v = max(int(v["verse"]) for v in ch_data["verses"])

            verses = get_verses(book, ch, start_v, end_v)
            texts.extend(verses)
        return texts
    else:
        verses_text = get_verses(book, chapter_start, verse_start, verse_end)
        return verses_text



def delete_all_justin_ppt():
    ppt_files = glob.glob("justin_*.pptx")
    for file in ppt_files:
        try:
            os.remove(file)
            print(f"'{file}' 파일 삭제 완료.")
        except PermissionError:
            print(f"'{file}' 파일이 열려 있어서 삭제하지 못했어")
        except Exception as e:
            print(f"파일 삭제 중 오류 발생: {e}")


# 1. 본문에 사용할 글꼴 이름 (폰트)
FONT_NAME = 'HY헤드라인M'  # 프레젠테이션 텍스트에 사용할 글꼴 이름

# 2. 글자 색상 (RGB 형식)
FONT_COLOR = RGBColor(255, 255, 255)  # 흰색 글자. RGB값: (Red, Green, Blue)

# 3. 글자 크기 (포인트 단위)
FONT_SIZE_TITLE_PT = 42  # 타이틀 28pt로 설정
FONT_SIZE_PT = 38  # 본문 28pt로 설정

# 4. 줄간격 비율
LINE_SPACING_FACTOR = 1.4  # 줄간격: 글자 크기 대비 1.2배 정도의 간격

# 5. 줄간격 (pt 단위에서 cm로 변환)
LINE_SPACING_PT = FONT_SIZE_PT * LINE_SPACING_FACTOR  # 줄간격 pt값 계산
LINE_HEIGHT_CM = LINE_SPACING_PT * 0.0352778  # 줄 하나가 차지하는 높이(cm). 1pt = 약 0.0352778cm

# 6. 단락 간 간격 (줄간격 외의 여유 간격)
SPACE_AFTER_PT = 10  # 단락 아래 여백(pt)
SPACE_AFTER_CM = SPACE_AFTER_PT * 0.0352778  # cm 단위로 변환

# 7. 본문 텍스트박스가 사용할 최대 높이 (슬라이드 안에서 넘치지 않도록 하기 위함)
MAX_BODY_HEIGHT_CM = 14  # 본문이 들어갈 수 있는 최대 높이(cm)

# 8. 한 줄에 들어갈 것으로 예상되는 최대 문자 수
CHARS_PER_LINE = 20  # 평균적으로 한 줄에 30자 정도 들어갈 것으로 가정하고 줄 수 계산

# 9. 슬라이드 배경 색
SLIDE_BG_COLOR = RGBColor(0, 0, 0)  # 검은색 배경


# === PPT 생성 함수 ===
def create_ppt_for_verse(prs, ref, texts, one_verse_per_slide=False):
    parsed = parse_reference(ref)
    if not parsed:
        print("Invalid reference:", ref)
        return

    book_name, chapter_start, verse_start, verse_end = parsed

    i = 0
    first_slide = True
    current_chapter = None

    # 선택한 비율에 따라 슬라이드 크기 설정
    if use_wide_slide.get():
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
    while i < len(texts):
        # 슬라이드 생성 및 배경 설정
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = SLIDE_BG_COLOR

        # 제목 텍스트박스 생성
        if prs.slide_width == Inches(13.33):
            title_shape = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(31.867), Cm(2))
        else:
            title_shape = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(23.3), Cm(2))
        title_tf = title_shape.text_frame
        title_tf.clear()
        p = title_tf.paragraphs[0]
        run = p.add_run()

        chap_num, _, _ = texts[i]
        full_book_name = full_names.get(book_name, book_name)

        if first_slide or chap_num != current_chapter:
            current_chapter = chap_num
            title_text = f"{full_book_name} {current_chapter}{'편' if full_book_name == '시편' else '장'}"
            body_y = Cm(2.5)
        else:
            title_text = ""
            body_y = Cm(0.5)

        run.text = title_text
        run.font.size = Pt(FONT_SIZE_TITLE_PT)
        run.font.name = FONT_NAME
        run.font.color.rgb = FONT_COLOR
        p.alignment = PP_ALIGN.LEFT

        # 본문 텍스트박스 생성
        if prs.slide_width == Inches(13.33):
            body_shape = slide.shapes.add_textbox(Cm(1), body_y, Cm(31.867), Cm(MAX_BODY_HEIGHT_CM))
        else:
            body_shape = slide.shapes.add_textbox(Cm(1), body_y, Cm(23.3), Cm(MAX_BODY_HEIGHT_CM))
        body_tf = body_shape.text_frame
        body_tf.word_wrap = True
        body_tf.clear()

        used_height = 0.0

        if one_verse_per_slide:
            # 절마다 한 슬라이드
            chap_num, verse_num, verse_text = texts[i]
            full_text = f"{verse_num}. {verse_text}"
            p = body_tf.add_paragraph()
            run = p.add_run()
            run.text = full_text
            run.font.size = Pt(FONT_SIZE_PT)
            run.font.name = FONT_NAME
            run.font.color.rgb = FONT_COLOR
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = Pt(LINE_SPACING_PT)
            p.space_after = Pt(SPACE_AFTER_PT)

            i += 1
        else:
            # 여러 절 한 슬라이드
            while i < len(texts):
                chap_num, verse_num, verse_text = texts[i]

                if chap_num != current_chapter:
                    break

                full_text = f"{verse_num}. {verse_text}"

                est_lines = math.ceil(len(full_text) / CHARS_PER_LINE)
                content_height = est_lines * LINE_HEIGHT_CM
                total_est_height = content_height + SPACE_AFTER_CM

                if used_height + total_est_height > MAX_BODY_HEIGHT_CM:
                    break

                p = body_tf.add_paragraph()
                run = p.add_run()
                run.text = full_text
                run.font.size = Pt(FONT_SIZE_PT)
                run.font.name = FONT_NAME
                run.font.color.rgb = FONT_COLOR
                p.alignment = PP_ALIGN.LEFT
                p.line_spacing = Pt(LINE_SPACING_PT)
                p.space_after = Pt(SPACE_AFTER_PT)

                used_height += total_est_height
                i += 1

        first_slide = False


        
# 프로그램 시작 시 삭제
delete_all_justin_ppt()



# 직접 키보드 입력받기 & PPT 생성 연동
print("예수전도단 홍천지부")
print()
print()

import subprocess
import platform

def open_file(filename):
    if platform.system() == "Windows":
        os.startfile(filename)
    elif platform.system() == "Darwin":  # macOS
        subprocess.call(("open", filename))
    else:  # Linux 등
        subprocess.call(("xdg-open", filename))

def clear_console():
    os.system('cls')

import threading
import subprocess
import tkinter as tk
import tkinter.ttk as ttk
from tkinter import messagebox

def generate_ppt_from_gui_input():
    # 진행 메시지 창 생성
    progress_win = tk.Toplevel(root)
    progress_win.title("진행중")
    progress_win.geometry("400x100")
    progress_win.resizable(False, False)
    tk.Label(progress_win, text="생성 중이야... 잠시만 기다려 줄래?", font=("맑은 고딕", 12)).pack(expand=True, padx=10, pady=20)

    progress_win.grab_set()  # 메시지창에 포커스 강제
    root.update()  # 메시지창이 바로 뜨도록 업데이트

    input_str = entry.get().strip()
    if input_str.lower() == 'exit':
        delete_all_justin_ppt()
        root.quit()
        return

    refs = [ref.strip() for ref in re.split(r'[,\s]+', input_str) if ref.strip()]
    if not refs:
        messagebox.showwarning("입력 오류", "성경 구절을 입력해야해..")
        progress_win.destroy()
        return

    prs = Presentation()
    valid_count = 0

    # 체크박스 변수 읽기 (예: one_verse_var)
    one_verse_per_slide = one_verse_var.get()

    for ref in refs:
        texts = get_texts_from_reference(ref)
        if texts:
            create_ppt_for_verse(prs, ref, texts, one_verse_per_slide=one_verse_per_slide)
            valid_count += 1
        else:
            messagebox.showerror("구절 오류", f"'{ref}' 구절을 찾을 수 없는걸..")
            progress_win.destroy()
            return

    if valid_count > 0:
        delete_all_justin_ppt()
        filename = f"justin_combined_{int(time.time())}.pptx"
        try:
            if os.path.exists(filename):
                os.remove(filename)
        except PermissionError:
            messagebox.showerror("파일 열림 오류", f"'{filename}' 파일이 열려있습니다. 닫고 다시 시도하세요.")
            progress_win.destroy()
            return

        progress_win.destroy()
        prs.save(filename)
        open_file(filename)
        # messagebox.showinfo("성공", f"'{filename}' 파일이 생성되었습니다.")


def show_bible_abbreviations():
    bible_order = [
        "창", "출", "레", "민", "신", "수", "삿", "룻", "삼상", "삼하",
        "왕상", "왕하", "대상", "대하", "스", "느", "에", "욥", "시",
        "잠", "전", "아", "사", "렘", "애", "겔", "단", "호", "욜",
        "암", "옵", "욘", "미", "나", "합", "습", "학", "슥", "말",
        "마", "막", "눅", "요", "행", "롬", "고전", "고후", "갈", "엡",
        "빌", "골", "살전", "살후", "딤전", "딤후", "딛", "몬", "히", "약",
        "벧전", "벧후", "요일", "요이", "요삼", "유", "계"
    ]


    root = tk.Tk()
    root.title("성경 줄임말 홍천버전")
    root.geometry("460x840")
    root.resizable(False, False)

    frame = ttk.Frame(root, padding=10)
    frame.pack(fill="both", expand=True)

    label = ttk.Label(frame, text="친절한 홍천지부", font=("맑은 고딕", 18, "bold"))
    label.pack(pady=(0, 15))

    text_frame = ttk.Frame(frame)
    text_frame.pack(fill="both", expand=True)

    scrollbar = ttk.Scrollbar(text_frame, orient="vertical")
    scrollbar.pack(side="right", fill="y")

    text_widget = tk.Text(text_frame, yscrollcommand=scrollbar.set, font=("맑은 고딕", 12))
    text_widget.pack(side="left", fill="both", expand=True)
    scrollbar.config(command=text_widget.yview)

    for abbr in bible_order:
        full = full_names.get(abbr, "정보없음")
        text_widget.insert("end", f"{abbr:<6} : {full}\n")

    text_widget.config(state="disabled")


import webbrowser
def open_homepage():
    webbrowser.open("https://www.hongcheondts.com/")  # 홍천DTS 홈페이지 주소 입력


def apply_font_settings():
    global FONT_SIZE_TITLE_PT, FONT_SIZE_PT, LINE_SPACING_FACTOR
    global LINE_SPACING_PT, LINE_HEIGHT_CM  # 이 두 변수도 업데이트해야 함

    try:
        FONT_SIZE_TITLE_PT = int(entry_title_size.get())
        FONT_SIZE_PT = int(entry_body_size.get())
        LINE_SPACING_FACTOR = float(entry_line_spacing.get())

        # 줄간격 관련 계산도 같이 다시 해줌
        LINE_SPACING_PT = FONT_SIZE_PT * LINE_SPACING_FACTOR
        LINE_HEIGHT_CM = LINE_SPACING_PT * 0.0352778

        save_settings()  # ← 적용 후 저장!
        
        messagebox.showinfo("적용 완료", "글자 크기 및 줄간격 설정이 적용되었습니다.")
    except ValueError:
        messagebox.showerror("입력 오류", "숫자만 입력해주세요.")

def save_settings():
    settings = {
        "FONT_SIZE_TITLE_PT": FONT_SIZE_TITLE_PT,
        "FONT_SIZE_PT": FONT_SIZE_PT,
        "LINE_SPACING_FACTOR": LINE_SPACING_FACTOR
    }
    with open("settings.json", "w", encoding="utf-8") as f:
        json.dump(settings, f)

def load_settings():
    global FONT_SIZE_TITLE_PT, FONT_SIZE_PT, LINE_SPACING_FACTOR
    global LINE_SPACING_PT, LINE_HEIGHT_CM

    try:
        with open("settings.json", "r", encoding="utf-8") as f:
            settings = json.load(f)
            FONT_SIZE_TITLE_PT = settings.get("FONT_SIZE_TITLE_PT", 46)
            FONT_SIZE_PT = settings.get("FONT_SIZE_PT", 42)
            LINE_SPACING_FACTOR = settings.get("LINE_SPACING_FACTOR", 1.4)

            # 줄간격 다시 계산
            LINE_SPACING_PT = FONT_SIZE_PT * LINE_SPACING_FACTOR
            LINE_HEIGHT_CM = LINE_SPACING_PT * 0.0352778
    except FileNotFoundError:
        pass  # 처음 실행 시 파일 없을 수 있음

load_settings()

root = tk.Tk()
root.title("예수전도단 홍천지부 성경ppt 생성기V1")
root.geometry("400x500")
root.resizable(False, False)

# 전체 패딩 주는 프레임
main_frame = tk.Frame(root, padx=30, pady=30)
main_frame.pack(expand=True)

# 안내 레이블
label = tk.Label(main_frame, text="성경 구절 입력 (예: 창1:1-5, 요3:16, 시23):", font=("맑은 고딕", 11))
label.pack(fill='x', pady=(0,10))

# 입력창
entry = tk.Entry(main_frame, width=50, font=("맑은 고딕", 12))
entry.pack(fill='x', pady=(0,15))

# PPT 생성 버튼
btn = ttk.Button(main_frame, text="PPT 생성할까?", command=generate_ppt_from_gui_input)
btn.pack(fill='x', pady=(0,10))

# 성경 줄임말 보기 버튼
abbr_btn = ttk.Button(main_frame, text="성경 줄임말 볼까?", command=show_bible_abbreviations)
abbr_btn.pack(fill='x', pady=(0,10))

#홍천홈피
homepage_btn = ttk.Button(main_frame, text="예수전도단 홍천지부 홈페이지 볼래?", command=open_homepage)
homepage_btn.pack(pady=5)

# 체크박스 프레임 (텍스트와 체크박스를 같이 좌측 정렬)
checkbox_frame = tk.Frame(main_frame)
checkbox_frame.pack(fill='x', pady=(0,15))

# 한 절씩 한 슬라이드
one_verse_var = tk.BooleanVar(value=True)
checkbox_one_verse = tk.Checkbutton(
    checkbox_frame,
    text="한 절씩 한 슬라이드",
    variable=one_verse_var,
    font=("맑은 고딕", 11)
)
checkbox_one_verse.pack(anchor='w')

# 와이드 슬라이드 여부
use_wide_slide = tk.BooleanVar(value=True)
checkbox_wide = tk.Checkbutton(
    checkbox_frame,
    text="와이드 (16:9)",
    variable=use_wide_slide,
    font=("맑은 고딕", 11)
)
checkbox_wide.pack(anchor='w')

# 글자 설정 프레임 생성 (main_frame 내부)
font_frame = tk.LabelFrame(main_frame, text="글자 크기 및 줄간격 설정", padx=10, pady=10)
font_frame.pack(fill='x', pady=10)

# 타이틀 글자 크기
tk.Label(font_frame, text="타이틀 글자 크기 (pt):").grid(row=0, column=0, sticky='e', padx=5, pady=5)
entry_title_size = tk.Entry(font_frame, width=10)
entry_title_size.grid(row=0, column=1, padx=5, pady=5)
entry_title_size.insert(0, str(FONT_SIZE_TITLE_PT))

# 본문 글자 크기
tk.Label(font_frame, text="본문 글자 크기 (pt):").grid(row=1, column=0, sticky='e', padx=5, pady=5)
entry_body_size = tk.Entry(font_frame, width=10)
entry_body_size.grid(row=1, column=1, padx=5, pady=5)
entry_body_size.insert(0, str(FONT_SIZE_PT))

# 줄간격 비율
tk.Label(font_frame, text="줄간격 비율 (예: 1.4):").grid(row=2, column=0, sticky='e', padx=5, pady=5)
entry_line_spacing = tk.Entry(font_frame, width=10)
entry_line_spacing.grid(row=2, column=1, padx=5, pady=5)
entry_line_spacing.insert(0, str(LINE_SPACING_FACTOR))

# 적용 버튼
apply_button = tk.Button(font_frame, text="적용", command=apply_font_settings)
apply_button.grid(row=3, column=0, columnspan=2, pady=10)



# 엔터키 이벤트 바인딩
root.bind('<Return>', lambda event: generate_ppt_from_gui_input())

root.mainloop()

